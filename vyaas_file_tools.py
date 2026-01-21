"""
VYAAS AI - File Creation Tools
Create HTML, Excel, Word, PowerPoint, PDF files and save to Desktop
"""

import os
import logging
from datetime import datetime
from livekit.agents import function_tool

logger = logging.getLogger("vyaas_file_tools")
logger.setLevel(logging.INFO)

# Get user's Desktop path
import termux_compatibility as termux

# Get user's Desktop path
if termux.is_android():
    # On Android Termux, map Desktop to Downloads folder for easy access
    DESKTOP_PATH = "/sdcard/Download"
    if not os.path.exists(DESKTOP_PATH):
        # Fallback to internal storage home if permission issue
        DESKTOP_PATH = os.path.expanduser("~")
else:
    DESKTOP_PATH = os.path.join(os.path.expanduser("~"), "Desktop")


def get_safe_filename(name: str, extension: str) -> str:
    """Generate a safe filename with timestamp to avoid overwriting."""
    # Remove invalid characters
    safe_name = "".join(c for c in name if c.isalnum() or c in (' ', '-', '_')).strip()
    if not safe_name:
        safe_name = "file"
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    return f"{safe_name}_{timestamp}.{extension}"


@function_tool()
async def create_html_file(filename: str, content: str) -> str:
    """
    Create an HTML file and save it to the Desktop.
    Use this to create web pages, calculators, games, or any HTML-based applications.
    
    Args:
        filename: Name for the file (without .html extension)
        content: Complete HTML content including <!DOCTYPE html>, <html>, <head>, <body> tags
                 Include inline CSS and JavaScript for best results.
    
    Returns:
        Success message with file path or error
    
    Example for a calculator:
        content should be a complete HTML document with:
        - <!DOCTYPE html>
        - <style> for CSS styling
        - <script> for JavaScript logic
        - Complete calculator UI with buttons
    """
    logger.info(f"Creating HTML file: {filename}")
    try:
        safe_filename = get_safe_filename(filename, "html")
        filepath = os.path.join(DESKTOP_PATH, safe_filename)
        
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(content)
        
        logger.info(f"HTML file created: {filepath}")
        return f"Done! HTML file saved to Desktop: {safe_filename}"
    except Exception as e:
        logger.error(f"HTML creation error: {e}")
        return f"Error creating HTML file: {str(e)}"


@function_tool()
async def create_text_file(filename: str, content: str) -> str:
    """
    Create a text file (.txt) and save it to the Desktop.
    
    Args:
        filename: Name for the file (without .txt extension)
        content: Text content to write to the file
    
    Returns:
        Success message with file path or error
    """
    logger.info(f"Creating text file: {filename}")
    try:
        safe_filename = get_safe_filename(filename, "txt")
        filepath = os.path.join(DESKTOP_PATH, safe_filename)
        
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(content)
        
        return f"Done! Text file saved to Desktop: {safe_filename}"
    except Exception as e:
        logger.error(f"Text file creation error: {e}")
        return f"Error creating text file: {str(e)}"


@function_tool()
async def create_excel_file(filename: str, data: str, headers: str = "") -> str:
    """
    Create an Excel (.xlsx) file and save it to the Desktop.
    
    Args:
        filename: Name for the file (without .xlsx extension)
        headers: Comma-separated column headers (e.g., "Name,Age,City")
        data: Rows of data, each row on a new line, columns separated by commas.
              Example: "John,25,Delhi
                        Priya,22,Mumbai
                        Rahul,28,Pune"
    
    Returns:
        Success message with file path or error
    """
    logger.info(f"Creating Excel file: {filename}")
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment
        
        wb = Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        
        # Add headers if provided
        row_num = 1
        if headers:
            header_cells = [h.strip() for h in headers.split(',')]
            for col, header in enumerate(header_cells, 1):
                cell = ws.cell(row=1, column=col, value=header)
                cell.font = Font(bold=True)
                cell.fill = PatternFill(start_color="4F81BD", end_color="4F81BD", fill_type="solid")
                cell.font = Font(bold=True, color="FFFFFF")
            row_num = 2
        
        # Add data rows
        if data:
            for row_data in data.strip().split('\n'):
                if row_data.strip():
                    cells = [cell.strip() for cell in row_data.split(',')]
                    for col, value in enumerate(cells, 1):
                        # Try to convert to number if possible
                        try:
                            value = float(value) if '.' in value else int(value)
                        except ValueError:
                            pass
                        ws.cell(row=row_num, column=col, value=value)
                    row_num += 1
        
        # Auto-adjust column widths
        for column in ws.columns:
            max_length = 0
            column_letter = column[0].column_letter
            for cell in column:
                try:
                    if len(str(cell.value)) > max_length:
                        max_length = len(str(cell.value))
                except:
                    pass
            ws.column_dimensions[column_letter].width = max_length + 2
        
        safe_filename = get_safe_filename(filename, "xlsx")
        filepath = os.path.join(DESKTOP_PATH, safe_filename)
        wb.save(filepath)
        
        return f"Done! Excel file saved to Desktop: {safe_filename}"
    except ImportError:
        return "Error: openpyxl library not installed. Run: pip install openpyxl"
    except Exception as e:
        logger.error(f"Excel creation error: {e}")
        return f"Error creating Excel file: {str(e)}"


@function_tool()
async def create_word_document(filename: str, title: str, content: str) -> str:
    """
    Create a Word document (.docx) and save it to the Desktop.
    
    Args:
        filename: Name for the file (without .docx extension)
        title: Document title (will be formatted as heading)
        content: Document content. Use \\n for new lines/paragraphs.
                 Use "## " prefix for subheadings (e.g., "## Introduction")
    
    Returns:
        Success message with file path or error
    """
    logger.info(f"Creating Word document: {filename}")
    try:
        from docx import Document
        from docx.shared import Pt, Inches
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        
        doc = Document()
        
        # Add title
        title_para = doc.add_heading(title, 0)
        title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Add content
        paragraphs = content.split('\n')
        for para in paragraphs:
            para = para.strip()
            if not para:
                continue
            
            if para.startswith('## '):
                # Subheading
                doc.add_heading(para[3:], level=1)
            elif para.startswith('### '):
                # Sub-subheading
                doc.add_heading(para[4:], level=2)
            else:
                # Normal paragraph
                p = doc.add_paragraph(para)
        
        safe_filename = get_safe_filename(filename, "docx")
        filepath = os.path.join(DESKTOP_PATH, safe_filename)
        doc.save(filepath)
        
        return f"Done! Word document saved to Desktop: {safe_filename}"
    except ImportError:
        return "Error: python-docx library not installed. Run: pip install python-docx"
    except Exception as e:
        logger.error(f"Word creation error: {e}")
        return f"Error creating Word document: {str(e)}"


@function_tool()
async def create_powerpoint(filename: str, title: str, slides: str) -> str:
    """
    Create a PowerPoint presentation (.pptx) and save it to the Desktop.
    
    Args:
        filename: Name for the file (without .pptx extension)
        title: Presentation title (for first slide)
        slides: Slide content in format:
                "Slide Title 1|Bullet 1|Bullet 2|Bullet 3
                 Slide Title 2|Bullet 1|Bullet 2"
                Each line is a slide. First item is title, rest are bullets.
                Separate items with | character.
    
    Returns:
        Success message with file path or error
    """
    logger.info(f"Creating PowerPoint: {filename}")
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        
        prs = Presentation()
        
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]  # Title slide
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = f"Created by Vyaas AI"
        
        # Add content slides
        bullet_slide_layout = prs.slide_layouts[1]  # Title and Content
        
        for slide_data in slides.strip().split('\n'):
            if not slide_data.strip():
                continue
            
            items = [item.strip() for item in slide_data.split('|')]
            if not items:
                continue
            
            slide = prs.slides.add_slide(bullet_slide_layout)
            slide.shapes.title.text = items[0]
            
            # Add bullet points
            if len(items) > 1:
                body_shape = slide.shapes.placeholders[1]
                tf = body_shape.text_frame
                tf.text = items[1] if len(items) > 1 else ""
                
                for bullet in items[2:]:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.level = 0
        
        safe_filename = get_safe_filename(filename, "pptx")
        filepath = os.path.join(DESKTOP_PATH, safe_filename)
        prs.save(filepath)
        
        return f"Done! PowerPoint saved to Desktop: {safe_filename}"
    except ImportError:
        return "Error: python-pptx library not installed. Run: pip install python-pptx"
    except Exception as e:
        logger.error(f"PowerPoint creation error: {e}")
        return f"Error creating PowerPoint: {str(e)}"


@function_tool()
async def create_pdf_document(filename: str, title: str, content: str) -> str:
    """
    Create a PDF document and save it to the Desktop.
    
    Args:
        filename: Name for the file (without .pdf extension)
        title: Document title
        content: Text content for the PDF. Use \\n for new lines.
    
    Returns:
        Success message with file path or error
    """
    logger.info(f"Creating PDF: {filename}")
    try:
        from fpdf import FPDF
        
        pdf = FPDF()
        pdf.add_page()
        
        # Set font - using built-in font for compatibility
        pdf.set_font("Arial", 'B', 24)
        pdf.cell(0, 20, title, 0, 1, 'C')
        
        pdf.set_font("Arial", '', 12)
        pdf.ln(10)
        
        # Add content
        for line in content.split('\n'):
            # Handle long lines by wrapping
            pdf.multi_cell(0, 8, line)
        
        safe_filename = get_safe_filename(filename, "pdf")
        filepath = os.path.join(DESKTOP_PATH, safe_filename)
        pdf.output(filepath)
        
        return f"Done! PDF saved to Desktop: {safe_filename}"
    except ImportError:
        return "Error: fpdf library not installed. Run: pip install fpdf"
    except Exception as e:
        logger.error(f"PDF creation error: {e}")
        return f"Error creating PDF: {str(e)}"


@function_tool()
async def create_python_file(filename: str, code: str) -> str:
    """
    Create a Python (.py) file and save it to the Desktop.
    
    Args:
        filename: Name for the file (without .py extension)
        code: Python code to write to the file
    
    Returns:
        Success message with file path or error
    """
    logger.info(f"Creating Python file: {filename}")
    try:
        safe_filename = get_safe_filename(filename, "py")
        filepath = os.path.join(DESKTOP_PATH, safe_filename)
        
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(code)
        
        return f"Done! Python file saved to Desktop: {safe_filename}"
    except Exception as e:
        logger.error(f"Python file creation error: {e}")
        return f"Error creating Python file: {str(e)}"


@function_tool()
async def list_desktop_files() -> str:
    """
    List all files on the user's Desktop.
    
    Returns:
        List of files on Desktop
    """
    try:
        files = os.listdir(DESKTOP_PATH)
        if not files:
            return "Desktop is empty"
        
        # Categorize files
        docs = []
        images = []
        others = []
        
        for f in files:
            ext = f.split('.')[-1].lower() if '.' in f else ''
            if ext in ['docx', 'doc', 'pdf', 'xlsx', 'xls', 'pptx', 'ppt', 'txt', 'html']:
                docs.append(f)
            elif ext in ['jpg', 'jpeg', 'png', 'gif', 'bmp']:
                images.append(f)
            else:
                others.append(f)
        
        result = f"Files on Desktop ({len(files)} total):\n"
        if docs:
            result += f"\n📄 Documents: {', '.join(docs[:10])}"
        if images:
            result += f"\n🖼️ Images: {', '.join(images[:10])}"
        if others:
            result += f"\n📁 Others: {', '.join(others[:10])}"
        
        return result
    except Exception as e:
        return f"Error listing files: {str(e)}"
